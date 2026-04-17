#!/usr/bin/env python3
"""
docs/00.input/ 파일 전처리 스크립트

지원 형식:
  - .xlsx, .xls  → openpyxl로 시트/셀 텍스트 추출
  - .pptx        → python-pptx로 슬라이드 텍스트 추출
  - .csv         → CP949/UTF-8 인코딩 자동 감지 후 파싱
  - .txt, .md    → 인코딩 자동 감지 후 읽기

Claude가 직접 읽는 형식 (추출 불필요):
  - .pdf, .png, .jpg, .jpeg, .gif, .webp

추출 결과: docs/00.input/extracted/*.md
"""

from pathlib import Path
from typing import Optional
import csv as csv_module
import sys

INPUT_DIR = Path(__file__).parent
EXTRACTED_DIR = INPUT_DIR / "extracted"

# 이 스크립트 자체와 숨김 파일 제외
SKIP_FILES = {"extract.py"}
# Claude가 직접 읽을 수 있는 형식 — 추출 불필요
NATIVE_EXTENSIONS = {".pdf", ".png", ".jpg", ".jpeg", ".gif", ".webp"}


# ── 인코딩 감지 ──────────────────────────────────────────────

def detect_encoding(path: Path) -> str:
    """chardet으로 인코딩 감지. 미설치 시 cp949 반환."""
    try:
        import chardet
        raw = path.read_bytes()
        result = chardet.detect(raw)
        return result.get("encoding") or "cp949"
    except ImportError:
        return "cp949"


def read_text_safe(path: Path) -> str:
    """UTF-8 → CP949 → chardet 감지 순으로 시도."""
    for encoding in ["utf-8-sig", "utf-8", "cp949", detect_encoding(path)]:
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    # 최후 수단: 깨진 문자는 ?로 대체
    return path.read_text(encoding="utf-8", errors="replace")


# ── 파일 타입별 추출 ─────────────────────────────────────────

def extract_excel(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        return f"[오류] openpyxl 미설치 — pip install openpyxl\n파일: {path.name}"

    wb = openpyxl.load_workbook(path, data_only=True)
    lines = [f"# {path.name}"]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"\n## 시트: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return f"[오류] python-pptx 미설치 — pip install python-pptx\n파일: {path.name}"

    prs = Presentation(path)
    lines = [f"# {path.name}"]
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## 슬라이드 {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def extract_csv(path: Path) -> str:
    content = read_text_safe(path)
    lines = [f"# {path.name}\n"]
    for row in csv_module.reader(content.splitlines()):
        if any(cell.strip() for cell in row):
            lines.append(" | ".join(row))
    return "\n".join(lines)


def extract_text(path: Path) -> str:
    return read_text_safe(path)


# ── 메인 처리 ────────────────────────────────────────────────

def process_file(path: Path) -> Optional[str]:
    """파일 타입에 맞는 추출 함수 호출. 처리 불필요 시 None 반환."""
    suffix = path.suffix.lower()

    if suffix in (".xlsx", ".xls"):
        return extract_excel(path)
    elif suffix == ".pptx":
        return extract_pptx(path)
    elif suffix == ".csv":
        return extract_csv(path)
    elif suffix in (".txt", ".md"):
        return extract_text(path)
    else:
        return None  # NATIVE_EXTENSIONS 또는 미지원 형식


def main():
    EXTRACTED_DIR.mkdir(exist_ok=True)

    processed = []
    native = []
    unsupported = []

    for file in sorted(INPUT_DIR.iterdir()):
        # 숨김 파일, 디렉터리, 스킵 목록 제외
        if file.name in SKIP_FILES or file.name.startswith(".") or file.is_dir():
            continue

        suffix = file.suffix.lower()

        if suffix in NATIVE_EXTENSIONS:
            native.append(file.name)
            continue

        result = process_file(file)

        if result is None:
            unsupported.append(file.name)
            continue

        out_path = EXTRACTED_DIR / (file.stem + ".md")
        out_path.write_text(result, encoding="utf-8")
        processed.append(file.name)
        print(f"✅ {file.name} → extracted/{file.stem}.md")

    if native:
        print(f"\n⏭️  Claude 직접 읽기 ({len(native)}개): {', '.join(native)}")
    if unsupported:
        print(f"⚠️  미지원 형식 ({len(unsupported)}개): {', '.join(unsupported)}")

    print(f"\n완료: 추출 {len(processed)}개 / 직접읽기 {len(native)}개 / 미지원 {len(unsupported)}개")


if __name__ == "__main__":
    main()
